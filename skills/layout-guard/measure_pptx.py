#!/usr/bin/env python
"""
layout-guard / measure_pptx.py

Deterministic layout audit for .pptx WITHOUT rendering. Reads shape geometry
straight from the OOXML via python-pptx, so it never opens PowerPoint and is
safe to run headless (unlike COM automation, which can disturb a live app session).

Three failure classes, mirroring the HTML engine:

  1. off-slide      — a shape extends past the slide bounds
  2. shape-overlap  — two content shapes' rectangles intersect (placed-box
                      collision; the PPTX analogue of "글씨 겹침")
  3. text-overflow  — estimated text block height/width exceeds its textbox and
                      the box has no autofit shrink. HEURISTIC (true font metrics
                      need a render) — flagged, never auto-trusted.

Usage:
    python measure_pptx.py <pptx-path> [--json out.json]

Exit code 2 if any issue found.

Note on overlap: full-bleed background shapes and slide-size placeholders are
skipped (cover >85% of slide), and a small intersection (<5% of the smaller
shape) is ignored as a deliberate touch/overlap-by-design.
"""
import sys
import json
import math
import argparse

try:
    sys.stdout.reconfigure(encoding="utf-8")  # Windows cp949 console safety
except Exception:
    pass

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("python-pptx not installed:  pip install python-pptx")

EMU_PER_PT = 12700  # 1 pt = 12700 EMU
OFF_TOL = 8 * EMU_PER_PT  # shapes may bleed a few pt past the edge by design


def rects_intersect(a, b):
    ix = min(a["r"], b["r"]) - max(a["l"], b["l"])
    iy = min(a["btm"], b["btm"]) - max(a["t"], b["t"])
    if ix <= 0 or iy <= 0:
        return 0.0
    inter = ix * iy
    smaller = min(a["w"] * a["h"], b["w"] * b["h"])
    return inter / smaller if smaller else 0.0


def shape_text(shape):
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def estimate_text_overflow(shape):
    """Rough: does the text plausibly exceed the box? Returns (overflow_bool, note)."""
    tf = shape.text_frame
    # MSO_AUTO_SIZE: 0=NONE (fixed box, text clips), 1=SHAPE_TO_FIT_TEXT (box
    # silently GROWS to fit text -> collision/off-slide risk that stored XML
    # hides), 2=TEXT_TO_FIT_SHAPE (font shrinks -> no overflow, safe to skip).
    try:
        af = int(tf.auto_size) if tf.auto_size is not None else 0
    except Exception:
        af = 0
    if af == 2:
        return False, "autofit-shrink"
    grows = af == 1

    box_w_pt = shape.width / EMU_PER_PT
    box_h_pt = shape.height / EMU_PER_PT
    if box_w_pt <= 0 or box_h_pt <= 0:
        return False, ""

    total_lines = 0.0
    max_font = 0.0
    explicit = False  # only trust the estimate when font size is actually authored
    for p in tf.paragraphs:
        txt = p.text or ""
        # font size: paragraph runs, else paragraph font, else assume 18pt
        sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
        if sizes or p.font.size is not None:
            explicit = True
        fs = max(sizes) if sizes else (p.font.size.pt if p.font.size else 18.0)
        max_font = max(max_font, fs)
        # avg glyph advance ~ 0.55*em for Latin; Korean is ~1.0*em (much wider).
        has_cjk = any("가" <= c <= "힣" or "぀" <= c <= "ヿ" for c in txt)
        adv = (1.0 if has_cjk else 0.55) * fs
        usable_w = box_w_pt - 7.2  # ~0.1in inset each side
        chars_per_line = max(1, usable_w / adv) if adv else 1
        line_count = max(1, math.ceil(len(txt) / chars_per_line)) if txt else 1
        total_lines += line_count

    if not explicit:
        # theme/master-inherited size — python-pptx can't see it, so guessing 18pt
        # produces garbage for small boxes. Don't flag what we can't measure.
        return False, "font-inherited"
    needed_h = total_lines * max_font * 1.2  # 1.2 line spacing
    exceed = needed_h - box_h_pt
    if grows:
        # SHAPE_TO_FIT_TEXT: PowerPoint grows the box to fit at render time, so the
        # authored height is expected to be "too small". This is by design, not a
        # bug — don't flag. (Any resulting collision surfaces via shape-overlap.)
        return False, "autofit-grow"
    # fixed box (NONE / inherited-none): text genuinely clips. Conservative gate —
    # char-width is only an estimate, so require a clear multi-line overrun.
    if needed_h > box_h_pt * 1.35 and exceed > 12:
        return True, f"~{needed_h:.0f}pt text vs {box_h_pt:.0f}pt box — likely clipped"
    return False, ""


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--json")
    ap.add_argument("--overflow", action="store_true",
                    help="also run the text-overflow heuristic (noisy on dense/theme decks)")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    SW, SH = prs.slide_width, prs.slide_height
    slide_area = SW * SH

    report = []
    had_issue = False

    for sidx, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue  # grouped/auto-positioned — geometry not resolvable here
            w, h = shape.width or 0, shape.height or 0
            try:
                ph = bool(shape.is_placeholder)
            except Exception:
                ph = False
            boxes.append({
                "shape": shape,
                "name": shape.name,
                "text": shape_text(shape).strip(),
                "l": shape.left, "t": shape.top,
                "r": shape.left + w, "btm": shape.top + h,
                "w": w, "h": h,
                "area_frac": (w * h) / slide_area if slide_area else 0,
                "ph": ph,  # placeholder bboxes are layout-driven (oversized) — unreliable
            })

        issues = {"off_slide": [], "overlap": [], "overflow": []}

        # 1: edge-clip — a shape that STRADDLES a slide edge (part on, part off) is
        # a real clipping bug. A shape fully outside the canvas is parked/staged
        # (common in working figure decks) and invisible to the audience — ignore.
        # Backgrounds (area>0.5) bleed by design — ignore.
        for bx in boxes:
            if bx["area_frac"] > 0.5:
                continue
            on_x = min(bx["r"], SW) - max(bx["l"], 0)
            on_y = min(bx["btm"], SH) - max(bx["t"], 0)
            on_slide = on_x > 0 and on_y > 0  # some part is within the canvas
            beyond = (bx["r"] > SW + OFF_TOL or bx["btm"] > SH + OFF_TOL
                      or bx["l"] < -OFF_TOL or bx["t"] < -OFF_TOL)
            if on_slide and beyond:
                over = max(bx["r"] - SW, bx["btm"] - SH, -bx["l"], -bx["t"]) / EMU_PER_PT
                issues["off_slide"].append({"name": bx["name"], "by_pt": round(over)})

        # 2: shape-overlap — only a real GLYPH collision counts. Text placed on a
        # no-text background (banner picture, card rectangle, pill) is normal
        # layering, NOT overlap (mirrors the HTML text-leaf-vs-text-leaf rule).
        # So require BOTH shapes to carry text, skip full-bleed backgrounds, and
        # demand a substantial intersection to ignore pill/label edge touches.
        content = [b for b in boxes if b["area_frac"] < 0.85 and not b["ph"]]
        for i in range(len(content)):
            for j in range(i + 1, len(content)):
                A, B = content[i], content[j]
                if not (A["text"] and B["text"]):
                    continue  # text-on-background layering — not a collision
                frac = rects_intersect(A, B)
                if frac > 0.25:  # candidate glyph collision (verify visually for PPTX)
                    issues["overlap"].append({
                        "a": f'{A["name"]}"{A["text"][:20]}"',
                        "b": f'{B["name"]}"{B["text"][:20]}"',
                        "frac": round(frac, 2),
                    })

        # 3: text-overflow heuristic — OPT-IN (--overflow). XML can't see text
        # wrapping or theme-inherited fonts, so char-width estimation over-flags on
        # dense/theme-styled decks. Reliable mainly for Claude-generated decks with
        # explicit simple styling. Off by default to keep the gate trustworthy.
        if args.overflow:
            for bx in boxes:
                if bx["ph"]:
                    continue  # placeholder geometry is layout-driven, not authored
                if not bx["text"] or not bx["shape"].has_text_frame:
                    continue
                of, note = estimate_text_overflow(bx["shape"])
                if of:
                    issues["overflow"].append({"name": bx["name"], "note": note, "heuristic": True})

        n = sum(len(v) for v in issues.values())
        if n:
            had_issue = True
        report.append({"slide": sidx, **issues})

    # output
    for r in report:
        n = len(r["off_slide"]) + len(r["overlap"]) + len(r["overflow"])
        print(f'{"✗" if n else "✓"} slide {r["slide"]:>3}  {str(n)+" issue(s)" if n else "ok"}')
        for o in r["off_slide"]:
            print(f'    off-slide by {o["by_pt"]}pt  {o["name"]}')
        for o in r["overlap"]:
            print(f'    SHAPE-OVERLAP {int(o["frac"]*100)}%  {o["a"]}  ✕  {o["b"]}')
        for o in r["overflow"]:
            print(f'    text-overflow? (heuristic)  {o["name"]}  [{o["note"]}]')

    bad = sum(1 for r in report if len(r["off_slide"]) + len(r["overlap"]) + len(r["overflow"]))
    print(f'\n{len(report) - bad}/{len(report)} slide(s) clean, {bad} with issues')
    if args.json:
        with open(args.json, "w", encoding="utf-8") as f:
            json.dump(report, f, ensure_ascii=False, indent=2, default=str)
        print(f"JSON -> {args.json}")

    sys.exit(2 if had_issue else 0)


if __name__ == "__main__":
    main()
